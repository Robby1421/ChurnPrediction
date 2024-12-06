import streamlit as st
from pptx import Presentation
import requests
from io import BytesIO

# Function to fetch the PowerPoint file from GitHub
def fetch_pptx_from_github(url):
    response = requests.get(url)
    if response.status_code == 200:
        return BytesIO(response.content)
    else:
        st.error("Failed to fetch the file from GitHub. Please check the URL.")
        return None

# Function to extract slides from a PowerPoint file
def extract_slides(file):
    presentation = Presentation(file)
    slides = []
    for slide in presentation.slides:
        content = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text
                content.append(text)
        slides.append("\n\n".join(content))
    return slides

# Streamlit app layout
st.title("Churn Prediction Use Case - GitHub Integration")

# GitHub file URL input
github_url = st.text_input("Enter the GitHub raw URL for the PowerPoint file:")
if github_url:
    pptx_file = fetch_pptx_from_github(github_url)
    if pptx_file:
        slides_content = extract_slides(pptx_file)
        total_slides = len(slides_content)

        if total_slides > 0:
            # Select the slide to display
            selected_slide = st.slider("Select Slide", 1, total_slides, 1)
            st.write(f"### Slide {selected_slide}:")
            st.write(slides_content[selected_slide - 1])
        else:
            st.warning("No slides were found in the fetched file.")
else:
    st.info("Enter the GitHub raw URL to begin.")

# Footer
st.markdown("---")
st.caption("Powered by Streamlit and python-pptx")
